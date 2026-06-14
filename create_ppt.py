
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1E, 0x27, 0x61)       # navy
BG_CARD = RGBColor(0x2A, 0x36, 0x7A)        # light navy
ACCENT = RGBColor(0x00, 0xA8, 0x96)          # teal
ACCENT2 = RGBColor(0xF9, 0x61, 0x67)         # coral
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCA, 0xDC, 0xFC)           # ice blue
GOLD = RGBColor(0xF9, 0xE7, 0x95)
GRAY = RGBColor(0x8A, 0x9D, 0xC0)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

# ======== SLIDE 1: Cover ========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# Accent bar
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)

# Title
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
    "Hermes Agent 小课程", font_size=44, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(8), Inches(1),
    "从零开始 21 节课，学会用 AI 助手帮你干活", font_size=22, color=LIGHT)

# Divider
add_shape(slide, Inches(1.5), Inches(4.5), Inches(3), Inches(0.06), ACCENT)

# Bottom info
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(8), Inches(0.6),
    "tccca1  |  Hermes Agent  |  2026", font_size=16, color=GRAY)

# ======== SLIDE 2: 课程概览 ========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, BG_DARK)

# Header
add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "课程体系", font_size=36, color=WHITE, bold=True)

# 5 category cards
categories = [
    ("🟢", "入门篇", "3 课", "好玩初体验\n搜索与总结\n读文件/改文件", ACCENT),
    ("🔵", "实用篇", "5 课", "写代码/跑脚本\nGitHub 一条龙\n搜索代码/重构\nPDF/文档处理\nPPT 演示文稿", ACCENT2),
    ("🟣", "创意篇", "5 课", "画架构图\n网页设计\n信息图\nAI 音乐\n手绘图表", GOLD),
    ("🟡", "进阶篇", "4 课", "定时任务\n博客监控\nGoogle 全家桶\nNotion 集成", RGBColor(0x02, 0x80, 0x90)),
    ("🟠", "硬核篇", "4 课", "本地大模型\nTDD 开发\n系统化调试\nClaude Code 代理", RGBColor(0xB8, 0x50, 0x42)),
]

for i, (icon, name, count, items, color) in enumerate(categories):
    left = Inches(0.5 + i * 2.5)
    top = Inches(1.8)
    card = add_shape(slide2, left, top, Inches(2.2), Inches(5.0), BG_CARD)
    card.shadow.inherit = False

    # Color top bar
    add_shape(slide2, left, top, Inches(2.2), Inches(0.08), color)

    add_text_box(slide2, Inches(left.inches + 0.2), Inches(top.inches + 0.3), Inches(1.8), Inches(0.5),
        f"{icon} {name}", font_size=20, color=color, bold=True)

    add_text_box(slide2, Inches(left.inches + 0.2), Inches(top.inches + 0.8), Inches(1.8), Inches(0.4),
        f"共 {count}", font_size=14, color=GRAY)

    # Item list
    lines = items.split("\n")
    for j, line in enumerate(lines):
        add_text_box(slide2, Inches(left.inches + 0.2), Inches(top.inches + 1.4 + j * 0.5), Inches(1.8), Inches(0.5),
            f"• {line}", font_size=13, color=LIGHT)

# ======== SLIDE 3: 已完成课程 ========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, BG_DARK)

add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "已学课程", font_size=36, color=WHITE, bold=True)

completed = [
    ("1", "ASCII 艺术", "在终端生成漂亮的文字艺术"),
    ("2", "搜索与总结", "搜网页、查资料、做摘要"),
    ("5", "GitHub 工作流", "创建仓库、提交代码、推送到远程"),
    ("6", "搜索代码/重构", "搜索函数定义、优化冗余逻辑"),
    ("7", "PDF 文档处理", "提取文字、拆分合并、转 Markdown"),
    ("9", "画架构图", "生成 SVG 架构图展示课程结构"),
]

for i, (num, title, desc) in enumerate(completed):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.0)
    top = Inches(1.8 + row * 1.8)

    card = add_shape(slide3, left, top, Inches(5.5), Inches(1.5), BG_CARD)

    # Number badge
    badge = add_shape(slide3, Inches(left.inches + 0.2), Inches(top.inches + 0.2), Inches(0.6), Inches(0.6), ACCENT, MSO_SHAPE.OVAL)
    badge.text_frame.paragraphs[0].text = num
    badge.text_frame.paragraphs[0].font.size = Pt(18)
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.paragraphs[0].space_before = Pt(0)
    badge.text_frame.paragraphs[0].space_after = Pt(0)
    badge.text_frame.word_wrap = False

    add_text_box(slide3, Inches(left.inches + 1.0), Inches(top.inches + 0.15), Inches(4.2), Inches(0.5),
        f"第 {num} 课: {title}", font_size=18, color=ACCENT, bold=True)
    add_text_box(slide3, Inches(left.inches + 1.0), Inches(top.inches + 0.7), Inches(4.2), Inches(0.6),
        desc, font_size=14, color=LIGHT)

# ======== SLIDE 4: 第8课实操演示 ========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, BG_DARK)

add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "第 8 课: PPT 演示文稿", font_size=36, color=WHITE, bold=True)

# Left: what was done
card_left = add_shape(slide4, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0), BG_CARD)
add_shape(slide4, Inches(0.8), Inches(1.8), Inches(0.08), Inches(5.0), ACCENT)

add_text_box(slide4, Inches(1.2), Inches(2.0), Inches(5.0), Inches(0.5),
    "用 python-pptx 从零创建", font_size=20, color=ACCENT, bold=True)

steps = [
    "1. 深色主题 (Navy #1E2761)",
    "2. 封面页: 标题 + 副标题 + 装饰条",
    "3. 课程概览: 5 大类别卡片布局",
    "4. 已完成课程: 带编号徽章的网格",
    "5. 本课实况: 左右分栏展示",
    "6. 文字 QA: markitdown 提取验证",
    "7. 视觉 QA: 转图片检查布局",
]
for i, step in enumerate(steps):
    add_text_box(slide4, Inches(1.2), Inches(2.7 + i * 0.55), Inches(5.0), Inches(0.5),
        step, font_size=14, color=LIGHT)

# Right: stats
card_right = add_shape(slide4, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
add_shape(slide4, Inches(7.0), Inches(1.8), Inches(0.08), Inches(5.0), ACCENT2)

add_text_box(slide4, Inches(7.4), Inches(2.0), Inches(4.5), Inches(0.5),
    "幻灯片详情", font_size=20, color=ACCENT2, bold=True)

stats = [
    ("4", "幻灯片"),
    ("python-pptx", "Python 库"),
    ("深色 + 卡片", "设计风格"),
    ("13.3×7.5\"", "宽屏尺寸"),
]
for i, (label, val) in enumerate(stats):
    add_text_box(slide4, Inches(7.4), Inches(2.7 + i * 1.0), Inches(4.5), Inches(0.5),
        val, font_size=28, color=WHITE, bold=True)
    add_text_box(slide4, Inches(7.4), Inches(3.3 + i * 1.0), Inches(4.5), Inches(0.4),
        label, font_size=14, color=GRAY)

prs.save(r"C:\Users\mi\projects\hermes-demo\Hermes课程介绍.pptx")
print("✅ PPT 已生成: Hermes课程介绍.pptx")
